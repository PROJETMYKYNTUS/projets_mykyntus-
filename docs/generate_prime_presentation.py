# -*- coding: utf-8 -*-
"""Génère la présentation PowerPoint Module PRIME (public non technique)."""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- Palette corporate ---
NAVY = RGBColor(0x1A, 0x2B, 0x4A)
TEAL = RGBColor(0x0D, 0x7C, 0x8C)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
SLATE = RGBColor(0x5A, 0x6B, 0x7D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
RED_SOFT = RGBColor(0xC4, 0x39, 0x39)
GREEN_SOFT = RGBColor(0x2E, 0x7D, 0x4A)

OUTPUT = Path(__file__).resolve().parent / "Module-PRIME-Presentation-Managers.pptx"
W, H = Inches(13.333), Inches(7.5)  # 16:9


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if hasattr(shape, "adjustments") and shape.adjustments:
        try:
            shape.adjustments[0] = 0.05
        except (IndexError, TypeError):
            pass
    return shape


def add_text_box(slide, left, top, width, height, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return box


def add_bullet_list(slide, left, top, width, height, items, size=16, color=NAVY, icon_ok=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "✔ " if icon_ok else "❌ "
        p.text = prefix + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return box


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_header(slide, title: str, subtitle: str | None = None):
    add_accent_bar(slide)
    add_text_box(slide, Inches(0.55), Inches(0.35), Inches(12), Inches(0.7), title, size=32, bold=True, color=NAVY)
    if subtitle:
        add_text_box(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.45), subtitle, size=14, color=SLATE)


def add_key_message(slide, text: str, top=Inches(6.35)):
    band = add_rect(slide, Inches(0.55), top, Inches(12.2), Inches(0.75), TEAL)
    add_text_box(slide, Inches(0.75), top + Inches(0.12), Inches(11.8), Inches(0.5), text, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_placeholder_frame(slide, left, top, width, height, label: str):
    frame = add_rect(slide, left, top, width, height, LIGHT_BG, TEAL)
    add_text_box(slide, left + Inches(0.15), top + height / 2 - Inches(0.35), width - Inches(0.3), Inches(0.7), label, size=13, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    inner = add_rect(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(1.1), WHITE, SLATE)
    add_text_box(slide, left + Inches(0.35), top + height / 2 - Inches(0.15), width - Inches(0.7), Inches(0.4), "INSÉRER CAPTURE ICI", size=11, color=SLATE, align=PP_ALIGN.CENTER)


def slide_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    # bande décorative
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), GOLD)
    add_rect(slide, Inches(8.5), Inches(0), Inches(4.833), H, RGBColor(0x14, 0x3D, 0x52))
    # formes abstraites business
    for i, (x, y, sz, col) in enumerate([
        (9.2, 1.2, 1.8, TEAL),
        (10.5, 3.5, 1.2, GOLD),
        (8.8, 5.2, 2.0, RGBColor(0x25, 0x5A, 0x6E)),
    ]):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(sz), Inches(sz))
        c.fill.solid()
        c.fill.fore_color.rgb = col
        c.line.fill.background()

    add_text_box(slide, Inches(0.75), Inches(1.8), Inches(7.5), Inches(1.1), "Module PRIME", size=44, bold=True, color=WHITE)
    add_text_box(
        slide, Inches(0.75), Inches(2.85), Inches(7.2), Inches(1.0),
        "Solution intelligente de gestion\ndes performances et des primes",
        size=20, color=RGBColor(0xD0, 0xE8, 0xEE),
    )
    add_rect(slide, Inches(0.75), Inches(4.0), Inches(1.8), Inches(0.06), GOLD)
    add_text_box(slide, Inches(0.75), Inches(4.35), Inches(3), Inches(0.5), "[ LOGO ENTREPRISE ]", size=12, color=SLATE)
    add_text_box(slide, Inches(0.75), Inches(5.5), Inches(5), Inches(0.35), "Nom entreprise", size=14, color=WHITE)
    add_text_box(slide, Inches(0.75), Inches(5.9), Inches(5), Inches(0.35), "Équipe / Direction", size=13, color=SLATE)
    add_text_box(slide, Inches(0.75), Inches(6.35), Inches(3), Inches(0.35), "Mai 2026", size=12, color=SLATE)
    add_text_box(slide, Inches(9.0), Inches(6.0), Inches(3.8), Inches(0.5), "Pilotage performance · Primes · Décision", size=11, color=RGBColor(0xA0, 0xC4, 0xCE), align=PP_ALIGN.RIGHT)


def slide_problem(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Pourquoi changer les méthodes classiques ?", "Les limites de la gestion manuelle")
    problems = [
        "Fichiers Excel complexes",
        "Temps perdu dans les calculs",
        "Risques d'erreurs humaines",
        "Difficulté de suivi des performances",
        "Manque de visibilité pour les managers",
        "Processus longs et répétitifs",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(1.55), Inches(5.5), Inches(4.2), problems, size=17, icon_ok=False)
    # Avant / Après visuel
    add_text_box(slide, Inches(6.4), Inches(1.55), Inches(2.5), Inches(0.4), "AVANT", size=14, bold=True, color=RED_SOFT, align=PP_ALIGN.CENTER)
    add_placeholder_frame(slide, Inches(6.2), Inches(2.0), Inches(2.9), Inches(2.0), "Excel · emails · fichiers dispersés")
    add_text_box(slide, Inches(9.6), Inches(1.55), Inches(2.5), Inches(0.4), "APRÈS", size=14, bold=True, color=GREEN_SOFT, align=PP_ALIGN.CENTER)
    add_placeholder_frame(slide, Inches(9.4), Inches(2.0), Inches(2.9), Inches(2.0), "Plateforme PRIME centralisée")
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.05), Inches(2.85), Inches(0.45), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEAL
    arrow.line.fill.background()
    add_key_message(slide, "Plus l'entreprise grandit, plus la gestion manuelle devient difficile.")


def slide_solution(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Une plateforme centralisée et intelligente", "Le Module PRIME en une vision")
    benefits = [
        "Centralisation des données",
        "Calcul automatique des primes",
        "Suivi des performances",
        "Dashboards clairs",
        "Gestion hiérarchique",
        "Visualisation rapide des KPI",
        "Historique et traçabilité",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(1.55), Inches(4.8), Inches(4.5), benefits, size=16)
    # Schéma hub central
    hub = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(2.6), Inches(2.2), Inches(1.1))
    hub.fill.solid()
    hub.fill.fore_color.rgb = TEAL
    hub.line.fill.background()
    add_text_box(slide, Inches(6.15), Inches(2.85), Inches(1.9), Inches(0.5), "PRIME", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    labels = [("Données", 5.0, 1.8), ("KPI", 8.5, 1.8), ("Primes", 5.0, 4.5), ("Managers", 8.5, 4.5)]
    for lbl, x, y in labels:
        n = add_rect(slide, Inches(x), Inches(y), Inches(1.5), Inches(0.65), LIGHT_BG, TEAL)
        add_text_box(slide, Inches(x), Inches(y + 0.15), Inches(1.5), Inches(0.4), lbl, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_placeholder_frame(slide, Inches(10.0), Inches(1.55), Inches(2.9), Inches(2.3), "Capture — Dashboard")
    add_placeholder_frame(slide, Inches(10.0), Inches(4.05), Inches(2.9), Inches(2.0), "Capture — Fiches primes")
    add_key_message(slide, "Une seule plateforme pour piloter efficacement les performances.")


def slide_business_value(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Valeur ajoutée concrète", "Les gains pour l'entreprise")
    items = [
        "Gain de temps important",
        "Réduction des erreurs humaines",
        "Meilleure productivité",
        "Suivi plus rapide des équipes",
        "Décisions plus rapides",
        "Données centralisées",
        "Organisation plus professionnelle",
        "Meilleure visibilité des performances",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(1.55), Inches(5.8), Inches(4.5), items, size=15)
    # KPI visuels illustratifs (pas des promesses chiffrées contractuelles)
    metrics = [("-70%", "tâches manuelles*"), ("-85%", "risque d'erreur*"), ("+40%", "réactivité décision*")]
    for i, (val, lbl) in enumerate(metrics):
        x = Inches(6.3 + i * 2.15)
        card = add_rect(slide, x, Inches(1.7), Inches(1.95), Inches(1.55), LIGHT_BG, TEAL)
        add_text_box(slide, x, Inches(1.95), Inches(1.95), Inches(0.7), val, size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(2.65), Inches(1.95), Inches(0.5), lbl, size=11, color=SLATE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.2), Inches(3.45), Inches(6.5), Inches(0.35), "* Ordres de grandeur illustratifs — à adapter avec vos données réelles", size=9, color=SLATE)
    # mini bar chart décoratif
    bars = [(0.35, 2.8), (0.55, 3.4), (0.75, 2.5), (0.95, 3.8), (1.15, 4.2)]
    for bx, bh in bars:
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5 + bx * 2), Inches(6.0 - bh * 0.35), Inches(0.35), Inches(bh * 0.35))
        b.fill.solid()
        b.fill.fore_color.rgb = TEAL if bh > 3 else RGBColor(0x6B, 0xB8, 0xC4)
        b.line.fill.background()
    add_text_box(slide, Inches(6.5), Inches(6.05), Inches(3), Inches(0.3), "Évolution productivité", size=10, color=SLATE)
    add_key_message(slide, "Moins de tâches manuelles, plus de pilotage intelligent.")


def slide_managers(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Une meilleure expérience managériale", "Plus de pilotage, moins d'administratif")
    items = [
        "Accès rapide aux KPI",
        "Suivi des équipes en temps réel",
        "Identification rapide des performances",
        "Historique clair et traçable",
        "Réduction des tâches administratives",
        "Meilleure prise de décision",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(1.55), Inches(5.2), Inches(4.2), items, size=16)
    add_placeholder_frame(slide, Inches(6.0), Inches(1.55), Inches(6.9), Inches(3.5), "Capture — Dashboard manager / KPI")
    # cartes KPI
    kpis = [("Taux validation", "87%"), ("Équipes suivies", "24"), ("Délai moyen", "-3 j")]
    for i, (t, v) in enumerate(kpis):
        x = Inches(6.2 + i * 2.25)
        card = add_rect(slide, x, Inches(5.25), Inches(2.05), Inches(0.95), NAVY)
        add_text_box(slide, x + Inches(0.1), Inches(5.35), Inches(1.85), Inches(0.35), t, size=10, color=RGBColor(0xA0, 0xC4, 0xCE))
        add_text_box(slide, x + Inches(0.1), Inches(5.65), Inches(1.85), Inches(0.45), v, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_key_message(slide, "Les managers passent moins de temps à chercher l'information et plus de temps à piloter.")


def slide_automation(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Automatiser avec plus de fiabilité", "Qualité et cohérence, pas seulement la vitesse")
    items = [
        "Réduction des erreurs Excel",
        "Standardisation des calculs",
        "Traitement plus rapide",
        "Cohérence des résultats",
        "Réduction des oublis manuels",
        "Génération rapide des fiches de primes",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(1.55), Inches(5.5), Inches(4.0), items, size=16)
    # flux simple
    steps = ["Saisie\nstructurée", "Calculs\nautomatisés", "Contrôles\nintégrés", "Fiches\nprimes"]
    for i, s in enumerate(steps):
        x = Inches(6.0 + i * 1.55)
        box = add_rect(slide, x, Inches(2.2), Inches(1.35), Inches(1.1), LIGHT_BG if i % 2 else TEAL, TEAL)
        col = NAVY if i % 2 else WHITE
        add_text_box(slide, x + Inches(0.05), Inches(2.45), Inches(1.25), Inches(0.8), s, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        if i < 3:
            arr = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x + Inches(1.38), Inches(2.55), Inches(0.2), Inches(0.35))
            arr.fill.solid()
            arr.fill.fore_color.rgb = GOLD
            arr.line.fill.background()
    add_placeholder_frame(slide, Inches(6.0), Inches(4.0), Inches(6.9), Inches(1.85), "Capture — Fiche de prime générée")
    add_key_message(slide, "L'objectif n'est pas seulement d'automatiser, mais d'améliorer la qualité du traitement.")


def slide_flexibility(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Une solution pensée pour le futur", "Flexibilité & évolutivité")
    items = [
        "Hiérarchie adaptable",
        "Ajout facile de nouveaux rôles",
        "Adaptation aux changements organisationnels",
        "Évolutivité long terme",
        "Possibilité d'ajouter de nouvelles fonctionnalités",
        "Solution scalable",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(1.55), Inches(5.0), Inches(3.8), items, size=15)
    add_text_box(slide, Inches(5.9), Inches(1.55), Inches(6.8), Inches(0.4), "Aujourd'hui", size=13, bold=True, color=TEAL)
    add_text_box(slide, Inches(5.9), Inches(1.95), Inches(6.8), Inches(0.5), "Manager → Superviseur → Coach → Pilote", size=14, color=NAVY)
    add_text_box(slide, Inches(5.9), Inches(2.65), Inches(6.8), Inches(0.4), "Demain (exemple)", size=13, bold=True, color=GOLD)
    add_text_box(
        slide, Inches(5.9), Inches(3.05), Inches(6.8), Inches(0.7),
        "Manager → Responsable Région → Superviseur → Coach → Pilote",
        size=13, color=NAVY,
    )
    # timeline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(4.2), Inches(6.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()
    for i, lbl in enumerate(["2026", "2027", "2028+"]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.2 + i * 3.0), Inches(4.05), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD if i == 2 else TEAL
        dot.line.fill.background()
        add_text_box(slide, Inches(5.9 + i * 3.0), Inches(4.35), Inches(1.0), Inches(0.35), lbl, size=11, color=SLATE, align=PP_ALIGN.CENTER)
    add_key_message(slide, "La solution peut évoluer avec l'entreprise sans refonte complète.")


def slide_vision(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Vision cible du Module PRIME", "Déjà en place · Prochaines étapes")
    add_text_box(slide, Inches(0.65), Inches(1.55), Inches(5.8), Inches(0.4), "✅ Déjà pris en charge aujourd'hui", size=15, bold=True, color=GREEN_SOFT)
    today = [
        "Traitement automatique des calculs",
        "Génération rapide des fiches de primes",
        "Suivi clair et traçable",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(2.0), Inches(5.5), Inches(2.0), today, size=14)
    add_text_box(slide, Inches(6.8), Inches(1.55), Inches(5.5), Inches(0.4), "Vision future", size=15, bold=True, color=TEAL)
    future = [
        "Import automatique multi-sources",
        "Centralisation complète des données",
        "Analyses intelligentes",
        "Dashboards avancés",
        "Reporting automatique",
        "Évolutions IA",
    ]
    add_bullet_list(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(3.5), future, size=14)
    add_rect(slide, Inches(0.65), Inches(4.5), Inches(12.0), Inches(1.35), LIGHT_BG, TEAL)
    add_text_box(slide, Inches(0.85), Inches(4.7), Inches(11.5), Inches(1.0), "Feuille de route : accompagner la croissance sans rupture opérationnelle", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_key_message(slide, "Le système a été conçu pour accompagner les besoins futurs de l'entreprise.")


def slide_demo(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Aperçu des interfaces", "Un produit prêt à l'emploi")
    frames = [
        ("Dashboard", 0.55, 1.45),
        ("Gestion KPI", 3.45, 1.45),
        ("Classements", 0.55, 4.05),
        ("Fiches de primes", 3.45, 4.05),
    ]
    for lbl, x, y in frames:
        add_placeholder_frame(slide, Inches(x), Inches(y), Inches(2.75), Inches(2.35), f"Capture — {lbl}")
    # deuxième rangée large
    add_placeholder_frame(slide, Inches(6.5), Inches(1.45), Inches(6.2), Inches(4.95), "Capture — Vue validation / workflow")
    add_text_box(slide, Inches(0.55), Inches(6.55), Inches(12), Inches(0.35), "Remplacez les zones grises par vos captures réelles avant la présentation.", size=10, color=SLATE, align=PP_ALIGN.CENTER)


def slide_conclusion(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), GOLD)
    add_text_box(slide, Inches(0.75), Inches(0.6), Inches(11), Inches(0.8), "Un outil pensé pour la performance", size=34, bold=True, color=WHITE)
    points = [
        "Gain de temps",
        "Réduction des erreurs",
        "Centralisation intelligente",
        "Meilleure visibilité",
        "Pilotage plus efficace",
        "Solution évolutive",
    ]
    cols = 2
    for i, p in enumerate(points):
        col = i % cols
        row = i // cols
        x = Inches(0.75 + col * 5.8)
        y = Inches(1.7 + row * 0.85)
        card = add_rect(slide, x, y, Inches(5.2), Inches(0.65), RGBColor(0x24, 0x3D, 0x5C), TEAL)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.12), Inches(4.8), Inches(0.45), "✔  " + p, size=16, color=WHITE)
    add_text_box(
        slide, Inches(0.75), Inches(4.5), Inches(11.5), Inches(1.2),
        "Le Module PRIME transforme la gestion des performances en un processus\nplus intelligent, plus fiable et plus efficace.",
        size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
    )


def slide_thanks(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)
    add_rect(slide, Inches(0), Inches(6.9), W, Inches(0.6), NAVY)
    add_text_box(slide, Inches(0.75), Inches(2.2), Inches(11.5), Inches(1.0), "Merci pour votre attention", size=40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.75), Inches(3.4), Inches(11.5), Inches(0.6), "Questions & échanges", size=22, color=TEAL, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.75), Inches(4.5), Inches(11.5), Inches(0.5), "[ Nom du présentateur ]  ·  [ Contact ]  ·  [ Entreprise ]", size=14, color=SLATE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.75), Inches(5.2), Inches(11.5), Inches(0.4), "Module PRIME — Pilotage des performances et des primes", size=12, color=SLATE, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_business_value(prs)
    slide_managers(prs)
    slide_automation(prs)
    slide_flexibility(prs)
    slide_vision(prs)
    slide_demo(prs)
    slide_conclusion(prs)
    slide_thanks(prs)
    prs.save(str(OUTPUT))
    print(f"Présentation créée : {OUTPUT}")
    print(f"Slides : {len(prs.slides)}")


if __name__ == "__main__":
    build()
